"""Build PPTX slide decks from structured content."""

import math
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string to an RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _calc_title_layout(title_text: str, avail_width_inches: float):
    """Pick title font size and box height so the text never overflows.

    Returns (font_size_pt, box_height_inches).
    """
    # Approximate chars-per-line at each font size for the available width.
    # At 45pt on ~5.7" width ≈ 22 cpl; 38pt ≈ 26 cpl; 32pt ≈ 31 cpl.
    tiers = [
        (45, avail_width_inches / 0.26),   # ~22 cpl at 5.7"
        (38, avail_width_inches / 0.22),   # ~26 cpl
        (32, avail_width_inches / 0.185),  # ~31 cpl
    ]

    max_lines = 4  # never exceed 4 lines regardless of font size
    line_height_factor = 1.35  # line height as multiple of font size

    for font_pt, cpl in tiers:
        lines = max(1, math.ceil(len(title_text) / max(cpl, 1)))
        if lines <= max_lines:
            line_h_in = (font_pt * line_height_factor) / 72  # pt to inches
            box_h = lines * line_h_in + 0.15  # small padding
            return font_pt, box_h

    # Fallback: smallest tier, cap at max_lines
    font_pt, cpl = tiers[-1]
    line_h_in = (font_pt * line_height_factor) / 72
    box_h = max_lines * line_h_in + 0.15
    return font_pt, box_h


def build_pptx(
    slides: list[dict],
    colors: dict,
    aspect_ratio: str = "9:16",
    output_dir: str = "./output",
    handle: str = "@cristian.bojaca",
) -> str:
    """Build a PPTX file from slide content.

    Args:
        slides: List of slide dicts with 'title', 'body'.
        colors: Dict with 'background', 'title', 'body', 'accent', 'highlight'.
        aspect_ratio: Either '9:16' (vertical/stories) or '16:9' (landscape).
        output_dir: Directory to save the output file.
        handle: Account handle shown at the bottom of each slide.

    Returns:
        Path to the generated PPTX file.
    """
    prs = Presentation()

    if aspect_ratio == "9:16":
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.33)
    else:
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    bg_color = _hex_to_rgb(colors.get("background", "#0D0D15"))
    title_color = _hex_to_rgb(colors.get("title", "#FFFFFF"))
    body_color = _hex_to_rgb(colors.get("body", "#C0C0D0"))
    accent_color = _hex_to_rgb(colors.get("accent", "#F7B731"))
    muted_color = _hex_to_rgb("#4A5568")

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.9)
    content_w = slide_w - 2 * margin
    content_w_inches = content_w / 914400  # EMU to inches

    accent_bar_w = Inches(0.05)

    for idx, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # -- Background --
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # -- Left accent bar: 5px, full height, accent blue --
        left_bar = slide.shapes.add_shape(
            1, 0, 0, accent_bar_w, slide_h,
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = accent_color
        left_bar.line.fill.background()

        # -- Slide counter: top-right, muted blue --
        indicator_w = Inches(1.5)
        indicator_box = slide.shapes.add_textbox(
            slide_w - margin - indicator_w,
            Inches(0.6),
            indicator_w,
            Inches(0.35),
        )
        indicator_tf = indicator_box.text_frame
        indicator_tf.word_wrap = True
        indicator_p = indicator_tf.paragraphs[0]
        indicator_p.text = f"{idx + 1}/{len(slides)}"
        indicator_p.font.size = Pt(16)
        indicator_p.font.name = "Arial"
        indicator_p.font.color.rgb = accent_color
        indicator_p.font.bold = False
        indicator_p.alignment = PP_ALIGN.RIGHT

        # -- Title: Arial Bold, dynamically sized --
        title_text = slide_data.get("title", "")
        title_font_pt, title_box_h = _calc_title_layout(title_text, content_w_inches)

        title_top = Inches(2.2)
        title_height = Inches(title_box_h)
        title_box = slide.shapes.add_textbox(
            margin, title_top, content_w, title_height
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.auto_size = None
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(title_font_pt)
        title_p.font.name = "Arial"
        title_p.font.bold = True
        title_p.font.color.rgb = title_color
        title_p.alignment = PP_ALIGN.LEFT
        title_p.space_after = Pt(0)

        # -- Thin horizontal divider (positioned after actual title height) --
        divider_top = title_top + title_height + Inches(0.15)
        divider = slide.shapes.add_shape(
            1,
            margin,
            divider_top,
            content_w,
            Inches(0.015),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = muted_color
        divider.line.fill.background()

        # -- Body: Times New Roman 35pt (positioned after divider) --
        body_top = divider_top + Inches(0.25)
        # Body fills remaining space above handle
        handle_zone = Inches(1.2)
        body_height = slide_h - body_top - handle_zone
        body_box = slide.shapes.add_textbox(
            margin, body_top, content_w, body_height
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.auto_size = None
        body_p = body_tf.paragraphs[0]
        body_p.text = slide_data.get("body", "")
        body_p.font.size = Pt(35)
        body_p.font.name = "Times New Roman"
        body_p.font.color.rgb = body_color
        body_p.alignment = PP_ALIGN.LEFT
        body_p.space_after = Pt(0)
        body_p.line_spacing = Pt(44)

        # -- Handle: centered, small, muted, bottom of slide --
        handle_box = slide.shapes.add_textbox(
            0,
            slide_h - Inches(0.9),
            slide_w,
            Inches(0.4),
        )
        handle_tf = handle_box.text_frame
        handle_tf.word_wrap = True
        handle_p = handle_tf.paragraphs[0]
        handle_p.text = handle
        handle_p.font.size = Pt(14)
        handle_p.font.name = "Arial"
        handle_p.font.color.rgb = muted_color
        handle_p.alignment = PP_ALIGN.CENTER

    # Save
    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"finance_slides_{timestamp}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return filepath
